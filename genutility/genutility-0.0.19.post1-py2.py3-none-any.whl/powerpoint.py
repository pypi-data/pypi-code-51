from pptx import Presentation # python-pptx

def read_ppt(path):
	# type: (str, ) -> Iterator[List[str]]

	prs = Presentation(path)
	for slide in prs.slides:
		shapes = [shape.text for shape in slide.shapes if hasattr(shape, "text")]
		yield shapes
